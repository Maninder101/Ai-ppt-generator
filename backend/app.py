from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import uuid
import re
import google.generativeai as genai
from dotenv import load_dotenv

# --------------------------
# Load Environment Variables
# --------------------------
load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")

if not GOOGLE_API_KEY:
    raise ValueError("❌ Missing GOOGLE_API_KEY in .env file!")

genai.configure(api_key=GOOGLE_API_KEY)
model = genai.GenerativeModel("gemini-2.5-flash")

# --------------------------
# Flask App Setup
# --------------------------
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
FRONTEND_BUILD = os.path.join(BASE_DIR, "../frontend/build")
OUTPUT_DIR = os.path.join(BASE_DIR, "generated")

app = Flask(__name__, static_folder=FRONTEND_BUILD, static_url_path="/")
CORS(app)

os.makedirs(OUTPUT_DIR, exist_ok=True)

# --------------------------
# Utility Functions
# --------------------------
def parse_slides(text, max_slides=6):
    slides = re.findall(r"Slide\s*\d+\s*:\s*(.*?)(?=Slide\s*\d+\s*:|$)", text, re.DOTALL)
    parsed = []
    for s in slides:
        lines = [l.strip(" -•\n") for l in s.strip().split("\n") if l.strip()]
        if lines:
            parsed.append({"title": lines[0], "points": lines[1:]})
    return parsed[:max_slides]

def apply_template_style(slide, template):
    theme_colors = {
        "modern": {"bg": RGBColor(25, 25, 112), "text": RGBColor(255, 255, 255)},
        "minimal": {"bg": RGBColor(255, 255, 255), "text": RGBColor(30, 30, 30)},
        "corporate": {"bg": RGBColor(0, 51, 102), "text": RGBColor(255, 215, 0)},
        "creative": {"bg": RGBColor(128, 0, 128), "text": RGBColor(255, 255, 255)},
        "dark": {"bg": RGBColor(15, 15, 15), "text": RGBColor(200, 200, 200)},
    }

    colors = theme_colors.get(template, theme_colors["modern"])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colors["bg"]
    return colors["text"]

def generate_ppt_content(topic, slide_count=6):
    prompt = f"""
    Create a PowerPoint outline on: "{topic}".
    Include exactly {slide_count} slides.
    Format:

    Slide 1: Title
    - Bullet 1
    - Bullet 2
    - Bullet 3
    """
    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        print("❌ Error:", e)
        return None

def create_ppt_from_slides(slides_data, template):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        text_color = apply_template_style(slide, template)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.3), Inches(1.2))
        title_tf = title_box.text_frame
        p = title_tf.add_paragraph()
        p.text = slide_data["title"]
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER

        # Points
        content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(10.3), Inches(4))
        content_tf = content_box.text_frame

        for bullet in slide_data["points"]:
            bp = content_tf.add_paragraph()
            bp.text = f"➤ {bullet.strip()}"
            bp.font.size = Pt(24)
            bp.font.color.rgb = text_color
            bp.level = 0

    filename = f"{uuid.uuid4()}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)
    return filename

# --------------------------
# API ROUTES
# --------------------------
@app.route('/generate-ppt', methods=['POST'])
def generate_ppt():
    data = request.get_json()
    topic_text = data.get('text', '').strip()
    template = data.get('template', 'modern')

    if not topic_text:
        return jsonify({"success": False, "error": "Empty topic"}), 400

    slides_data = parse_slides(topic_text)
    if not slides_data:
        return jsonify({"success": False, "error": "Invalid slide format"}), 400

    filename = create_ppt_from_slides(slides_data, template)
    return jsonify({
        "success": True,
        "download_url": f"/generated/{filename}"
    })

@app.route('/generate-auto-ppt', methods=['POST'])
def generate_auto_ppt():
    data = request.get_json()
    topic = data.get('topic', '').strip()
    slide_count = int(data.get("slide_count", 6))
    template = data.get("template", "modern")

    if not topic:
        return jsonify({"success": False, "error": "Topic missing"}), 400

    ai_text = generate_ppt_content(topic, slide_count)
    if not ai_text:
        return jsonify({"success": False, "error": "AI failed"}), 500

    slides_data = parse_slides(ai_text, slide_count)
    filename = create_ppt_from_slides(slides_data, template)

    return jsonify({
        "success": True,
        "slides_generated": len(slides_data),
        "download_url": f"/generated/{filename}"
    })

@app.route('/generated/<path:filename>')
def download_file(filename):
    return send_from_directory(OUTPUT_DIR, filename, as_attachment=True)

# --------------------------
# REACT FRONTEND SERVE
# --------------------------
@app.route('/')
def serve_react():
    return send_from_directory(FRONTEND_BUILD, "index.html")

@app.errorhandler(404)
def fallback(e):
    return send_from_directory(FRONTEND_BUILD, "index.html")

# --------------------------
# Run Server
# --------------------------
if __name__ == '__main__':
    app.run(host="0.0.0.0", port=5000)
